import json
from pptx import Presentation

from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

from pptx.util import Pt


def save_image(shape, slide_index, image_index, output_dir):
    # 获取图片的二进制数据
    image_stream = shape.image.blob
    # 构造图片文件名
    image_filename = f"slide_{slide_index}_image_{image_index}.png"
    image_path = os.path.join(output_dir, image_filename)
    # 将图片数据写入文件
    with open(image_path, 'wb') as img_file:
        img_file.write(image_stream)
    return image_filename

def read_ppt_to_json(ppt_file, output_dir='../../output/images'):
    # 确保输出目录存在
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # 载入PPT文件
    presentation = Presentation(ppt_file)

    # 初始化字典来保存幻灯片文本和图片信息
    slides_data = {}

    manu_script = []

    # 遍历幻灯片
    for i, slide in enumerate(presentation.slides):
        slides_data[i] = {'manuscript': [], 'images': []}
        image_count = 0

        # 遍历幻灯片中的形状
        for shape in slide.shapes:
            if shape.has_text_frame:

                text = ' '.join(paragraph.text for paragraph in shape.text_frame.paragraphs)
                if text.strip() != '':
                    # 识别文本类型：根据文本框位置或字体大小
                    font_size = shape.text_frame.paragraphs[0].runs[0].font.size

                    text_data = {'text': text, 'type': 'content'}
                    # 假设最大的字体大小是标题
                    if font_size and font_size > Pt(32):  # Pt是字体大小单位
                        text_data['type'] = 'title'
                        manu_script.append(f"# {text_data['text']}")
                    # 假设较大的字体大小是副标题
                    elif font_size and font_size > Pt(24):
                        text_data['type'] = 'subtitle'
                        manu_script.append(f"## {text_data['text']}")
                    else:
                        manu_script.append(f"### {text_data['text']}")

                    slides_data[i]['manuscript'].append(text_data)

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 保存图片并记录到JSON中
                image_filename = save_image(shape, i, image_count, output_dir)
                slides_data[i]['images'].append(image_filename)
                image_count += 1

    # 将字典转换为JSON字符串
    json_text = json.dumps(slides_data, indent=4, ensure_ascii=False)
    manu_script_text = '\n'.join(manu_script)
    return json_text, manu_script_text

# Use the function
ppt_file = '..\\..\\sample_reports\\电动化变革热度高涨，双轮车迎来出海新机会.pptx'  # replace with your ppt file path
json_output, manuscript = read_ppt_to_json(ppt_file)

# Save the output to a file
with open('..\\..\\output\\output.json', 'w', encoding='utf-8') as f:
    f.write(json_output)


with open('../../data/workflows/MarketPlan/input/raw_data/manuscript.txt', 'w', encoding='utf-8') as f:
    f.write(manuscript)